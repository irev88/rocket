# -*- coding: utf-8 -*-
"""Build Day 6 PPT deck — 可重复使用策略（回收概念）, fully Simplified Chinese."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

NAVY   = RGBColor(0x16, 0x32, 0x4F)
BLUE   = RGBColor(0x2E, 0x5E, 0x8C)
ORANGE = RGBColor(0xE0, 0x56, 0x1E)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GRAY   = RGBColor(0x5B, 0x66, 0x70)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF9)
CARD   = RGBColor(0xEE, 0xF3, 0xF8)
PEACH  = RGBColor(0xFB, 0xEE, 0xE6)
MINT   = RGBColor(0xE8, 0xF2, 0xEC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x33, 0x47, 0x5B)
FONT   = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
PAGE = [0]

def _set_font(run, size, bold=False, color=DARK, italic=False):
    f = run.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", FONT)

def rect(slide, x, y, w, h, fill, line=None, lw=1.0, shadow=False, round_=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, space_after=2):
    """lines: list of (text, size, bold, color) or list of lists of such runs per para."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after)
        runs = ln if isinstance(ln, list) else [ln]
        for (t, size, bold, color) in runs:
            r = p.add_run(); r.text = t; _set_font(r, size, bold, color)
    return tb

def bullet_block(slide, x, y, w, h, items, size=15, color=DARK, gap=6, mark_color=ORANGE):
    lines = []
    for it in items:
        if isinstance(it, tuple):  # (lead, rest)
            lines.append([("▪ ", size, True, mark_color), (it[0], size, True, NAVY), (it[1], size, False, color)])
        else:
            lines.append([("▪ ", size, True, mark_color), (it, size, False, color)])
    return text(slide, x, y, w, h, lines, space_after=gap)

def header(slide, title, sub=None):
    rect(slide, 0, 0, 13.333, 0.92, NAVY)
    rect(slide, 0, 0.92, 13.333, 0.045, ORANGE)
    text(slide, 0.55, 0.12, 11.5, 0.7, [(title, 27, True, WHITE)])
    if sub:
        text(slide, 9.2, 0.30, 3.6, 0.5, [(sub, 12, False, RGBColor(0xBF,0xD7,0xEA))], align=PP_ALIGN.RIGHT)
    PAGE[0] += 1
    rect(slide, 0, 7.18, 13.333, 0.32, LIGHT)
    text(slide, 0.55, 7.19, 9.5, 0.3,
         [("可重复使用运载火箭 AI 协同设计暑期课程 · Day 6 可重复使用策略（回收概念）· 2026-07-16", 10.5, False, GRAY)])
    text(slide, 12.2, 7.19, 0.7, 0.3, [(str(PAGE[0]), 10.5, False, GRAY)], align=PP_ALIGN.RIGHT)

def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s

def style_cell(cell, txt, size=12.5, bold=False, color=DARK, fill=None, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.MIDDLE):
    cell.margin_left = Pt(5); cell.margin_right = Pt(4)
    cell.margin_top = Pt(1.5); cell.margin_bottom = Pt(1.5)
    cell.vertical_anchor = anchor
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for i, (t, b, c) in enumerate(txt):
        r = p.add_run(); r.text = t; _set_font(r, size, b, c)

def make_table(slide, x, y, w, h, data, col_w=None, header_fill=NAVY, body_size=12.5,
               head_size=13, band=True, first_col_bold=False):
    rows, cols = len(data), len(data[0])
    gfx = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gfx.table
    tbl.first_row = False; tbl.horz_banding = False
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / total))
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                style_cell(cell, [(val, True, WHITE)], size=head_size, fill=header_fill,
                           align=PP_ALIGN.CENTER)
            else:
                fill = WHITE
                if band and ri % 2 == 0: fill = CARD
                bold = first_col_bold and ci == 0
                col = NAVY if bold else DARK
                if isinstance(val, tuple):
                    bold, col = val[1], val[2]
                    style_cell(cell, [(val[0], bold, col)], size=body_size, fill=fill)
                else:
                    style_cell(cell, [(val, bold, col)], size=body_size, fill=fill)
    return tbl

# ============================================================ SLIDE 1 — TITLE
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 4.62, 13.333, 0.05, ORANGE)
text(s, 0.9, 1.05, 11.5, 0.5, [("可重复使用运载火箭 AI 协同设计暑期课程（11–20 July 2026）", 15, False, RGBColor(0xBF,0xD7,0xEA))])
text(s, 0.9, 1.75, 11.6, 1.6, [("可重复使用策略", 48, True, WHITE), ("（回收概念 · Recovery Concept）", 26, False, RGBColor(0xE8,0xF2,0xEC))])
text(s, 0.9, 3.55, 11.5, 0.8, [("动力垂直回收 · 远海网/绳捕获 · 复用护照 —— 面向 20 t SSO 任务的一级复用方案", 17, False, RGBColor(0xBF,0xD7,0xEA))])
rect(s, 0.9, 4.95, 4.9, 1.5, RGBColor(0x1F,0x45,0x6E), round_=True)
text(s, 1.12, 5.12, 4.5, 1.2, [("Day 6 / 10 · 2026年7月16日", 16, True, WHITE),
                               ("交付物：回收概念（Recovery Concept）", 12.5, False, RGBColor(0xBF,0xD7,0xEA))])
rect(s, 6.0, 4.95, 6.4, 1.5, RGBColor(0x1F,0x45,0x6E), round_=True)
text(s, 6.22, 5.12, 6.1, 1.3, [("技术基线：Day 1–4 工作包成果（需求 / Δv 预算 / 发动机 / 质量与架构）", 12.5, False, RGBColor(0xBF,0xD7,0xEA)),
                               ("外部锚点：长征十号B（2026-07-10）· 星舰 F5（2024-10-13）· Block 5 运营经济", 12.5, False, RGBColor(0xBF,0xD7,0xEA))])
text(s, 0.9, 6.85, 11.6, 0.4, [("说明：本档量化基础不含 Day 5 仿真产品（Day 7 校核结论：模型已由 Day 7 重建并完成校验）；估计值均已在文中标注方法与量级锚点｜v1.1（2026-07-17 校勘增补）", 11.5, False, RGBColor(0x9A,0xA7,0xB2))])
notes(s, "Day 6 交付presentation。基线来自 Day 1–4；外部来源见报告参考文献 [1]–[23]。")

# ============================================================ SLIDE 2 — INPUTS
s = prs.slides.add_slide(BLANK); header(s, "任务承接：来自 Day 1–4 的约束输入", "Inputs")
data = [
 ["输入项", "采用值", "来源", "在本日工作中的作用"],
 ["任务剖面", "20,000 kg 至太阳同步轨道（SSO）", "Day 1", "高能、高倾角 ⇒ 远海回收走廊（第3–5章）"],
 ["成本与频次", "< $30M/次；≥ 10 次/年", "Day 1", "决定寿命、周转与舰队目标（第13页）"],
 ["回收 Δv 分配", "再入 500–800 / 末段 200–300 m/s", "Day 2 §2.3", "第12页储备闭合的燃烧配额"],
 ["发动机（Merlin 级）", "845 kN（海平面）×9+1；节流 40–100%；可多次重启", "Day 3", "末段点火物理（第10页）、再入点火（第9页）"],
 ["一级质量模型", "干重 40 t；回收储备 18 t；回收硬件 5.5 t", "Day 4 §5", "储备闭合（式6.1）与捕获能量（E∝v²）"],
 ["末端回收架构", "混合式海上捕获：捕获环＋4捕获耳＋磨损靴＋最小缓冲垫", "Day 4 §4.4–4.6", "绑定决策（第6、11页）；着陆腿仅作备份"],
 ["制导与传感", "惯性/GNSS → 走廊控制 → 船端相对导航；FBG 复用护照", "Day 4 §3.2/§8.2", "终端制导与复飞认证（第13页）"],
]
make_table(s, 0.55, 1.25, 12.25, 5.0, data, col_w=[2.0,4.6,1.5,4.1], body_size=12.5, head_size=13.5)
text(s, 0.55, 6.42, 12.2, 0.5, [[("依据声明：", 13, True, RED), ("本日量化基础仅取 Day 1–4 成果与外部引用运行数据；Day 5 OpenRocket 产品待 Day 7 轨迹模型修复后回填精确状态量。", 13, False, GRAY)]])
notes(s, "表：报告第1.2节 Table 1 的演示版。")

# ============================================================ SLIDE 3 — WHY REUSE
s = prs.slides.add_slide(BLANK); header(s, "为什么复用：一级是发射成本的主体", "Economics")
cards = [
 ("60–75%", "助推器占发射总成本的比例（2026 S-1 时代数据 [1]）"),
 ("~$30M ⇔ <$30万", "新造一枚助推器 vs 翻修一枚回收助推器（< 1%）[1][2]"),
 ("≈$2,500/kg", "复用后 LEO 单位成本：约 −75%（原 $9,000–10,000 级）[1]"),
 ("157 / 165", "2025 年 Falcon 9 任务中使用复飞助推器的次数（复用=默认）[2]"),
 ("25 / 40 / 36", "会计折旧寿命 / 工程设计目标 / 现役纪录（B1067 飞行次数）[1][2]"),
 ("第 2 次回本", "翻修 < 新造成本 10%；第 2 次飞行回本、第 3 次起净省 [4]"),
]
for i, (big, small) in enumerate(cards):
    cx = 0.55 + (i % 3) * 4.18; cy = 1.30 + (i // 3) * 1.85
    rect(s, cx, cy, 3.95, 1.62, CARD if i % 2 == 0 else PEACH, round_=True)
    text(s, cx + 0.2, cy + 0.14, 3.6, 0.6, [(big, 25, True, ORANGE if i % 2 else NAVY)])
    text(s, cx + 0.2, cy + 0.80, 3.6, 0.75, [(small, 12.5, False, DARK)])
rect(s, 0.55, 5.15, 12.25, 1.55, LIGHT, round_=True)
text(s, 0.8, 5.3, 12.0, 0.35, [("当前运营实践（锚定我们的检查与周转设计）", 14.5, True, NAVY)])
bullet_block(s, 0.85, 5.72, 11.8, 0.9, [
 ("翻修实践：", "逐台发动机检查测试；储箱焊缝与压力容器超声/X 射线检测微裂纹；预防性更换涡轮；必要时整机更换；复飞前全箭静态点火 [1]"),
 ("周转基准：", "现役平均约 40 天，最优 < 3 周 ⇒ 本项目目标：早期 ≤30 天、成熟 ≤14 天（第13页）"),
], size=12.5, gap=4)
notes(s, "来源：TechTimes 2026-06-12 / 2026-07-14（S-1 披露折旧 25 次、工程目标 40）；Wikipedia 复用开发计划（Musk 2020）。")

# ============================================================ SLIDE 4 — METHOD ZOO
s = prs.slides.add_slide(BLANK); header(s, "回收方法全景：六类技术路线对比", "Trade study")
data = [
 ["#", "方法家族", "代表案例（引用）", "成熟度", "优势", "劣势"],
 ["1", "降落伞＋海上溅落", "航天飞机 SRB；电子号海上回收 [7][8]", "成熟", "无需发动机重启", "精度低；海水腐蚀是主要翻修负担"],
 ["2", "降落伞＋空中捕获", "电子号直升机捕获（2022 短暂捕获后放弃）[7][9]", "演示后弃用", "避免海水接触", "操作脆弱；仅 ~50% 任务可用（vs 海面 60–70%）"],
 ["3", "动力垂直·着陆腿", "Falcon 9（600+ 次）；新格伦（2025-11 首落）[1][10]", "成熟运营", "轨道级验证最充分；海上/陆上两用", "腿质量＋触地载荷；展开失效模式"],
 ["4", "动力垂直·塔臂捕获", "星舰超重 F5（2024-10-13 首次即捕）[11][12]", "已演示", "取消腿质量；可“落地即复位”", "需返场推进剂；失效危及地面设施"],
 ["5", "动力垂直·网/绳捕获", "长征十号B（2026-07-10，世界首次）[13–16]", "首次演示", "取消腿质量；远海免返场；船端吸能", "捕获载荷集中；海况约束；公开数据有限"],
 ["6", "带翼水平起降", "航天飞机；X-37B", "轨道级历史", "大横向机动；回收环境温和", "翼/起落架质量；研制成本最高"],
]
make_table(s, 0.55, 1.25, 12.25, 5.5, data, col_w=[0.5,2.3,3.6,1.4,2.4,3.4], body_size=12, head_size=13, first_col_bold=True)
notes(s, "评分准则：① 与 Day 3 发动机能力匹配；② 固定 600 t GLOM 下的飞行质量代价；③ 20 t SSO 任务几何适配；④ 技术就绪度；⑤ 翻修负担；⑥ 基础设施风险。")

# ============================================================ SLIDE 5 — SELECTION
s = prs.slides.add_slide(BLANK); header(s, "选择结论：动力垂直回收 · 远海模式", "Selection")
bullet_block(s, 0.55, 1.25, 12.2, 2.3, [
 ("推进匹配：", "Day 3 已选发动机天然具备深度节流（40–100%）与多次重启能力——只有动力回收能利用这些既有能力（降落伞/带翼均无法利用）；"),
 ("任务几何：", "SSO 为高倾角高能任务，助推器天然远下程飞行，远海回收顺应弹道；"),
 ("外部验证：", "动力垂直家族已有 >600 次轨道级成功回收（腿式）＋ 2024 塔臂 ＋ 2026 网捕两次捕获演示 [1][10][11][13]；"),
 ("战略延展：", "不依赖跑道/机场，与月面无跑道场景兼容（课程长期目标的工程一致性）。"),
], size=14.5, gap=8)
rect(s, 0.55, 3.62, 12.25, 0.5, NAVY, round_=False)
text(s, 0.8, 3.68, 12.0, 0.35, [("决定性数字：回收模式的载荷处罚（SpaceX 公开数据 [5]）", 15, True, WHITE)])
rect(s, 0.55, 4.32, 12.25, 1.9, CARD, round_=True)
text(s, 1.0, 4.55, 3.2, 0.5, [("远海（船降）", 14, True, NAVY)])
rect(s, 1.0, 5.02, 3.1, 0.42, GREEN)
text(s, 4.25, 5.06, 2.5, 0.4, [("≈ 15%", 16, True, GREEN)])
text(s, 1.0, 5.62, 4.4, 0.5, [("返场（RTLS，含返场点火）", 14, True, NAVY)])
rect(s, 1.0, 6.02, 6.2, 0.42, RED)
text(s, 7.35, 6.06, 2.5, 0.4, [("≈ 30%", 16, True, RED)])
text(s, 8.1, 4.62, 4.5, 1.5, [("20,000 kg 载荷为 Day 1 不可谈判项 ——", 13.5, True, DARK),
                             ("双倍的性能处罚使返场在本任务级别被直接否决；远海为唯一可行模式。", 13.5, False, GRAY)])
notes(s, "CSMonitor 2016（Musk 2013 记者会原文）；Falcon Heavy 侧助推返场+芯级抛弃约 10% 作对照 [6]。")

# ============================================================ SLIDE 6 — TERMINAL ARCHITECTURE
s = prs.slides.add_slide(BLANK); header(s, "末端构型决策：混合式海上捕获（绑定 Day 4 结论）", "Terminal architecture")
data = [
 ["构型", "飞行硬件", "回收硬件质量", "回收储备", "主要风险", "结论"],
 ["着陆腿式", "着陆腿＋栅格翼＋RCS＋TPS＋传感器", "8,000 kg", "24,000 kg", "触地载荷；腿展开失效", ("成熟先例 → 保留为备份", True, GRAY)],
 ["纯网/绳捕获", "捕获钩/耳＋栅格翼＋RCS＋TPS", "< 5,500 kg（估）", "≈18,000 kg（估）", "单点挂接即全失", ("不可作为唯一系统", True, GRAY)],
 ["混合捕获兼容", "捕获环＋4捕获耳（磨损靴）＋栅格翼＋选择性TPS＋传感器＋最小缓冲垫", "5,500 kg", "18,000 kg", "捕获载荷设计；船端阻尼系统", ("选定基线", True, ORANGE)],
]
make_table(s, 0.55, 1.30, 12.25, 2.9, data, col_w=[1.7,4.6,1.6,1.4,2.2,2.1], body_size=12, head_size=13)
rect(s, 0.55, 4.45, 12.25, 2.25, MINT, round_=True)
text(s, 0.8, 4.62, 12.0, 0.4, [("在固定 GLOM = 600,000 kg 下，混合捕获相对着陆腿备份方案：", 14.5, True, NAVY)])
bullet_block(s, 0.85, 5.08, 11.8, 1.5, [
 ("释放质量：", "回收硬件 −2,500 kg（Day 7 校勘：储备推进剂需求由再入热物性决定，两者相同，不单独节省 −6 t）；"),
 ("再投资：", "+5,000 kg 上升推进剂、+4,500 kg 干重增长裕度（Day 4 §5.6 闭合）；"),
 ("预算一致性：", "远海合计低界 ≈0.8 km/s 的可实现性，正依赖船端承担末端支撑职能——这是捕获方案在预算层面的自我印证（第12页）。"),
], size=13.5, gap=5)
notes(s, "表摘自 Day 4 §4.4/§5.6 并经本日复核；kinetic energy 由船端液压阻尼吸收。")

# ============================================================ SLIDE 7 — EXTERNAL ANCHORS
s = prs.slides.add_slide(BLANK); header(s, "外部验证锚点：捕获路线已被飞行证明", "Flight heritage")
rect(s, 0.55, 1.25, 6.0, 4.15, CARD, round_=True)
text(s, 0.8, 1.42, 5.6, 0.4, [("长征十号B · 世界首次网捕（10 天前）", 16.5, True, ORANGE)])
bullet_block(s, 0.8, 1.95, 5.5, 3.3, [
 "2026-07-10 04:15 UTC，文昌 LC-2 首飞 [13–16]",
 "一级：7× YF-100K（液氧煤油，55 s 级燃烧后分离）",
 "三发动机再入点火 → 四栅格翼受控再入",
 "四只上置捕获钩挂入船端张紧钢索",
 "分离后约 6 分钟完成捕获；船位远海 ≈430 km（南海）",
 "“领航者号”：首个升高立体捕获架＋可动张紧钢索",
 "CASC 计划：同一枚助推器 2026 年内复飞",
], size=12.8, gap=4.5)
rect(s, 6.75, 1.25, 6.0, 4.15, CARD, round_=True)
text(s, 7.0, 1.42, 5.6, 0.4, [("星舰超重 F5/F7 · 塔臂捕获（2024–25）", 16.5, True, BLUE)])
bullet_block(s, 7.0, 1.95, 5.5, 3.3, [
 "F5（2024-10-13）：首次尝试即捕获成功 [11]",
 "捕获前数千项箭/塔状态判据逐项确认",
 "F7（2025-01-16）：一台发动机返场点火未复燃，捕获仍成功 [12]",
 "动机（Musk）：取消着陆腿 = 减重 + 取消腿成本",
 "“复位即复飞”：直接放回发射台，目标 1 小时内可复飞",
 "差异：塔捕需返场推进剂；本方案为远海模式，选网/绳路线",
], size=12.8, gap=4.5)
rect(s, 0.55, 5.55, 12.25, 1.15, PEACH, line=RED, lw=1.2, round_=True)
text(s, 0.8, 5.68, 12.0, 0.9, [
 [("警示证据：", 14, True, RED), ("朱雀三号（2025-12-03）与长征十二号A（2025-12-23）均抵达回收区，但双双失利于", 13, False, DARK), ("末段点火异常", 13, True, RED), (" [14]", 11, False, GRAY)],
 [("⇒ 末段点火的点火可靠性/行为位列本方案失效模式之首（F-1，第14页）。", 13, True, NAVY)],
], space_after=4)
notes(s, "NASASpaceFlight 2026-07-15 China update；Ars Technica 2024-10-13；19FortyFive 2026-07-12。")

# ============================================================ SLIDE 8 — SEQUENCE MAP
s = prs.slides.add_slide(BLANK); header(s, "回收飞行序列总览（七阶段）", "Sequence")
s.shapes.add_picture("assets/diagrams/svg_recovery_sequence.png", Inches(1.56), Inches(1.22), width=Inches(10.22))
notes(s, "图为矢量示意图；状态量为 Falcon 9 / LM-10B 级锚点估计，本箭精确值待 Day 7。")

# ============================================================ SLIDE 9 — PHASE NUMBERS
s = prs.slides.add_slide(BLANK); header(s, "阶段物理 ①：分离 → 惯性飞行 → 再入点火", "Phases 1–3")
data = [
 ["事件", "时间量级（估）", "状态量级（估）", "物理要点"],
 ["级间分离·翻转·沉降", "T+0–60 s（分离 T+142 s，Day 7 模型）", "≈66.5 km；1.89 km/s；γ≈41°（Day 7 模型）", "冷气 RCS 翻转 180°；挡板/防涡器沉降残余推进剂保复点燃供给"],
 ["惯性弹道飞行", "T+约 4–6 min", "顶点 ≈135 km（按分离态 v⊥≈1.23 km/s 复算）", "h_apo = h_sep + v⊥²/2g；水平漂移 ≈1.4 km/s 决定船位"],
 ["再入点火（3 台）", "T+约 5–7 min，约 15–30 s", "≈55→40 km；Δv 500–800 m/s", "3×~60% 节流 ⇒ T/W≈2.7；峰值热/力载荷封顶；F9/LM-10B 同为 3 台"],
 ["栅格翼气动下滑", "T+约 6–8 min", "高超声→亚声速；~250–300 m/s", "钛翼耐 >1,000℃ 免更换；走廊控制；转入船端相对导航"],
]
make_table(s, 0.55, 1.26, 12.25, 3.75, data, col_w=[2.6,2.8,2.9,5.6], body_size=11.5, head_size=12.5)
rect(s, 0.55, 5.22, 12.25, 1.60, LIGHT, round_=True)
bullet_block(s, 0.85, 5.36, 11.8, 1.4, [
 ("锚点：", "复用型 Falcon 9 分离 ≈Mach 6 级（vs 抛弃型 Mach 10 级）[17][18]；GTO 剖面分离 ≈8,350 km/h [19]；LM-10B 三发动机再入点火 [14]"),
 ("本箭模型值（Day 7 修复后）：", "分离 T+142 s / 66.5 km / 1,892 m/s / Mach 5.7 / γ=41° / 下程 51 km —— 回收 Δv 配额 0.8–1.2 km/s 与船位结论仍成立（敏感度：分离速度 −8%）"),
 ("校勘声明：", "本表数值锚点已替换为 Day 7 校验模型输出；剩余验证项 = 再入/末段子段细化（V-4 蒙特卡洛，Day 7-8）"),
], size=11.5, gap=3)
notes(s, "Day 7 模型输出已回填（原锚点估计）；v1.1 校勘。")

# ============================================================ SLIDE 10 — HOVER SLAM
s = prs.slides.add_slide(BLANK); header(s, "阶段物理 ②：末段点火与悬停猛击", "Terminal burn")
s.shapes.add_picture("assets/diagrams/svg_hover_slam.png", Inches(6.45), Inches(1.30), width=Inches(6.45))
bullet_block(s, 0.55, 1.35, 5.75, 4.6, [
 ("悬停判据：", "捕获前质量 ≈42–45 t ⇒ 重量 412–441 kN；单台 338–845 kN（40–100%）"), 
 ("悬停所需节流 ≈49–52%：", "恰落于不确定区间（文档下限 40% / 运行报道 ~50–57%）⇒ 不能依赖悬停"),
 ("基准机动 = 悬停猛击：", "在点火窗口内启动，速度 ≤2 m/s 与高度归零同步；只允许向上补节流、到顶即关机"),
 ("失败包容弹道：", "末段弹道偏置——发动机未点燃/欠推力时近船落海，不撞船"),
 ("捕获接口目标：", "垂直速度 ≤2 m/s · 横向偏移 ≤5 m · 姿态误差 ≤3° · 角速度 ≤1°/s（工程目标，待验证 V-2/V-4）"),
], size=13.5, gap=8)
notes(s, "物理：燃尽助推器在最低节流附近 T/W≈1，无法稳定悬停（r/nasa 等运行社区分析 [22]，与 Day 2 内部记载一致）。")

# ============================================================ SLIDE 11 — LOAD PATH
s = prs.slides.add_slide(BLANK); header(s, "捕获力学与结构负载路径", "Capture mechanics")
s.shapes.add_picture("assets/diagrams/svg_catch_loadpath.png", Inches(1.85), Inches(1.22), width=Inches(9.64))
notes(s, "船端吸收 E∝v²；磨损靴/TPS 为耗材；捕获带局部等网格加强。")

# ============================================================ SLIDE 12 — RESERVE CLOSURE
s = prs.slides.add_slide(BLANK); header(s, "回收推进剂预算闭合（18 t 储备）", "Δv closure")
s.shapes.add_picture("assets/diagrams/svg_reserve_closure.png", Inches(0.60), Inches(1.30), width=Inches(7.6))
bullet_block(s, 8.5, 1.55, 4.35, 4.3, [
 ("可用：", "Δv = Isp·g·ln(58/40 t) = 1.03–1.13 km/s（Isp 282–311 s）"),
 ("需求：", "再入 0.5–0.8 ＋ 末段 0.2–0.3 ＋ 余量 ~0.1 ⇒ 0.8–1.2 km/s"),
 ("结论 ①：", "远海模式闭合（上限偏紧）"),
 ("结论 ②：", "返场模式 1.1–1.7 km/s —— 不拨款；备份腿式 24 t 储备（≈1.38 km/s）亦仅勉强靠边 ⇒ 架构仅远海"),
 ("若需变更：", "触发 Day 4 §6.5 敏感性规则（Δm_上升推进剂 = −Δm_干重）"),
], size=13, gap=8)
rect(s, 8.5, 5.75, 4.35, 0.95, MINT, round_=True)
text(s, 8.72, 5.9, 3.95, 0.7, [("储备不仅够用，而且定义了任务形态：", 12.5, True, GREEN), ("无返场点火、无悬停徘徊。", 12.5, False, DARK)])
notes(s, "式（6.1）（6.2），报告第6节。")

# ============================================================ SLIDE 13 — OPERATIONS
s = prs.slides.add_slide(BLANK); header(s, "复飞运营：检查 · 翻修 · 寿命目标", "Operations")
steps = ["接收·安全处理", "复用护照下载", "定向检查(UT/X射线)", "发动机检查·积碳管理", "更换耗材(靴/TPS/垫)", "传感器校准·检漏", "静态点火 ⇒ 放行"]
cx = 0.35
for i, st in enumerate(steps):
    shp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(cx), Inches(1.30), Inches(1.86), Inches(0.85))
    shp.fill.solid(); shp.fill.fore_color.rgb = NAVY if i < 6 else GREEN
    shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = st; _set_font(r, 10.5, True, WHITE)
    cx += 1.83
data = [
 ["指标", "目标值", "外部锚点"],
 ["助推器设计寿命", "15 次（冲刺目标 40 次）", "Block 5 级：折旧 25 / 工程目标 40 / 纪录 36 [1][2][4]"],
 ["舰队频次", "≥ 10 次/年（Day 1 要求）", "165 次/年（2025 现役量级）"],
 ["周转时间", "早期 ≤30 天；成熟 ≤14 天", "现役均值 ~40 天、最优 <3 周 [2]"],
 ["翻修成本限额", "$5M/次（Day 1 模型）", "对照 ~$30 万级基准 [1] ⇒ Day 8 校准项 V-8"],
 ["现役助推器数", "2–3 枚轮换", "10×30/365 ≈ 0.8 ⇒ 含损耗裕度"],
]
make_table(s, 0.55, 2.55, 12.25, 3.0, data, col_w=[2.6,3.4,4.4], body_size=12.5, head_size=13, first_col_bold=True)
rect(s, 0.55, 5.75, 12.25, 1.0, CARD, round_=True)
text(s, 0.8, 5.88, 12.0, 0.75, [
 [("复用护照（Day 4 §3.2 设计特性）：", 13.5, True, NAVY), ("逐飞次记录应变/热/加速度历史、捕获载荷重构、发动机复点燃史——把“全面拆解”变为“定向检查”：超限区检查、其余放行。", 12.5, False, DARK)],
], space_after=2)
notes(s, "检查集合锚定现役实践：焊缝无损检测、涡轮预防更换、整机偶换、复飞前静态点火 [1][2]。")

# ============================================================ SLIDE 14 — RISKS
s = prs.slides.add_slide(BLANK); header(s, "失效模式登记与验证移交", "Risk & V&V")
data = [
 ["编号", "失效模式", "证据 / 理由", "设计缓解", "验证"],
 ["F-1", ("末段点火异常（点火/推力行为）", True, RED), "朱雀三号、长征十二号A（2025-12）均失利于此 [14]", "复点燃合格化；避船弹道；缓冲垫", "V-2"],
 ["F-2", "横向偏差超出捕获窗口", "风场、导航误差、船体运动", "船端相对导航；走廊控制；窗口设计", "V-4"],
 ["F-3", "捕获速度超限", "点火过晚/推力不足", "向上节流余量；阻尼容量；90→563 kJ 能量图", "V-2"],
 ["F-4", "非对称挂接（2–3 耳）", "捕获瞬间姿态/角速度误差", "4 耳几何；捕获环分配；FBG 见证", "V-5"],
 ["F-5", "捕获耳磨损/啃伤", "重复滑移接触", "可更换磨损靴；护照阈值", "V-6"],
 ["F-6/7", "壳体过载 · 沉降供给扰动", "集中捕获力；低液位晃动", "局部等网格；挡板/防涡器", "V-5/V-2"],
 ["F-8/9", "捕获失败落水 · 海况不可用", "非正常接触；气象窗口", "缓冲垫；海况 ≤4 假设（4 m 浪级船已验证 [14]）", "Day 8"],
]
make_table(s, 0.55, 1.22, 12.25, 4.35, data, col_w=[0.8,3.1,3.3,3.4,1.0], body_size=11.5, head_size=12.5, first_col_bold=True)
rect(s, 0.55, 5.75, 12.25, 1.0, LIGHT, round_=True)
text(s, 0.8, 5.86, 12.0, 0.85, [
 [("验证移交（报告表 8）：", 13, True, NAVY), ("V-1 — 已由 Day 7 模型提供初值（T+142 s · 66.5 km · 1,892 m/s · γ 41°）；V-2 点火窗口与节流下限判明 · V-3 悬吊稳定性 · V-4 散布蒙特卡洛 · V-5 捕获带结构筛查 · V-6 磨损靴寿命 ⇒ Day 7/8；V-7 海况可用性 · V-8 翻修成本校准 ⇒ Day 8", 11.5, False, DARK)],
], space_after=2)
notes(s, "完整 F-1…F-10 与 V-1…V-8 见 Day 6 报告第 8 节。")

# ============================================================ SLIDE 15 — VERDICT
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 1.55, 13.333, 0.05, ORANGE)
text(s, 0.7, 0.45, 12.0, 0.9, [("最终结论：远海 · 混合捕获 · 动力垂直回收 —— 回收侧闭合且被验证", 26, True, WHITE)])
items = [
 ("质量闭合", "相对着陆腿备份：回收硬件 −2.5 t（Day 7 校勘：储备推进剂由下行热物性驱动，两者相同；全箭统一锁定 ≥34.5 t）"),
 ("推进剂闭合", "18 t 储备 ⇒ 1.03–1.13 km/s 可用，对远海需求 0.8–1.2 km/s；返场模式被储备与 15%/30% 处罚双重否决"),
 ("需求可追溯", "贯穿 Day 1（20 t SSO / <$30M / ≥10 次/年）与 Day 3（深度节流、重启、复用耐受）的全部约束"),
 ("外部验证", "长征十号B 网捕（2026-07-10，直接先例，~6 min、远海 430 km）＋ 星舰 F5 塔捕（减重逻辑）＋ Block 5 运营经济（寿命与翻修基准）"),
 ("风险聚焦", "两次 2025-12 末段点火失利案例 → F-1 置顶；避船弹道与缓冲垫作内置安全特性；任务成功与回收成功解耦"),
]
y = 1.95
for t_, d_ in items:
    rect(s, 0.7, y, 0.16, 0.80, ORANGE)
    text(s, 1.05, y, 3.0, 0.8, [(t_, 16, True, RGBColor(0xF2,0xA3,0x3C))])
    text(s, 3.9, y + 0.02, 8.9, 0.9, [(d_, 13.5, False, WHITE)])
    y += 1.02
rect(s, 0.7, 6.28, 12.0, 1.0, RGBColor(0x1F,0x45,0x6E), round_=True)
text(s, 0.95, 6.38, 11.6, 0.85, [
 [("Day 7 校勘（2026-07-17）：", 12, True, RGBColor(0xF2,0xA3,0x3C)),
  ("600 t 文档化架构对 500×500 任务缺口 2,088 m/s（600 t 复用能力 ≈6.6 t；闭合 = 架构重标定 ≈802 t 或载荷调整 ⇒ Day 8 决策）——不影响本页回收侧结论。", 11.5, False, RGBColor(0xBF,0xD7,0xEA))],
 [("移交：", 12, True, RGBColor(0xF2,0xA3,0x3C)), ("V-1 已初值 · V-2…V-6 → Day 7（设计迭代）· V-7/V-8 → Day 8（可靠性与经济性）｜ v1.1 校勘版", 11.5, False, RGBColor(0xBF,0xD7,0xEA))],
], space_after=3)
notes(s, "详细论证与 23 项参考文献见《Day 6 可重复使用策略（回收概念）》最终版报告。")

prs.save("presentations/day06_reusability.pptx")
print("Saved presentations/day06_reusability.pptx, slides:", len(prs.slides._sldIdLst))
