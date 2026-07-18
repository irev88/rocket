#!/usr/bin/env python3
"""Heuristic layout audit for Day6 PPTX: flags probable text overflow and out-of-bounds shapes."""
from pptx import Presentation
from pptx.util import Emu

PPTX = "/home/user/Day6_可重复使用策略_回收概念.pptx"
prs = Presentation(PPTX)
SW = prs.slide_width
SH = prs.slide_height
PT_PER_EMU = 72.0 / 914400.0

def emu2pt(v):
    return v * PT_PER_EMU

def char_units(ch):
    o = ord(ch)
    if o >= 0x2E80:      # CJK & full-width
        return 1.0
    if ch in "iljtm.,:;!|()[] ":
        return 0.32
    if ch.isdigit():
        return 0.56
    if ch.isupper():
        return 0.68
    return 0.52

def para_lines(text, size_pt, width_pt):
    units = sum(char_units(c) for c in text)
    units_per_line = max(width_pt / size_pt, 1.0)
    return max(1, -(-int(units) // max(int(units_per_line), 1))) if units else 1

def frame_requirement(tf, width_pt):
    total = 0.0
    for p in tf.paragraphs:
        sizes = [r.font.size.pt for r in p.runs if r.font.size]
        size = max(sizes) if sizes else 18.0
        text = "".join(r.text for r in p.runs)
        ls = p.line_spacing if p.line_spacing else 1.18
        line_h = size * (ls if isinstance(ls, float) else 1.18)
        n = para_lines(text, size, width_pt)
        total += n * line_h
        if p.space_after:
            total += p.space_after.pt
        if p.space_before:
            total += p.space_before.pt
    # internal margins
    mi = tf.margin_top or 0
    mb = tf.margin_bottom or 0
    total += emu2pt(mi) + emu2pt(mb)
    return total

issues = []
for idx, slide in enumerate(prs.slides, 1):
    for shp in slide.shapes:
        name = shp.shape_type
        l, t, w, h = shp.left or 0, shp.top or 0, shp.width or 0, shp.height or 0
        # bounds check
        if l < -10000 or t < -10000 or l + w > SW + 10000 or t + h > SH + 10000:
            issues.append(f"S{idx}: OUT-OF-BOUNDS {shp.shape_id} '{(shp.name or '')[:24]}' L={emu2pt(l):.0f} T={emu2pt(t):.0f} R={emu2pt(l+w):.0f} B={emu2pt(t+h):.0f} (slide {emu2pt(SW):.0f}x{emu2pt(SH):.0f})")
        if shp.has_table:
            tbl = shp.table
            # estimate per-row required height from cell text
            declared = sum(r.height for r in tbl.rows)
            wpt = emu2pt(w)
            ncol = len(tbl.columns)
            # column widths
            colw = [emu2pt(c.width) for c in tbl.columns]
            need = 0.0
            for r in tbl.rows:
                row_need = 0.0
                for ci, cell in enumerate(r.cells):
                    size = 12.0
                    txt = cell.text
                    sizes = []
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.size:
                                sizes.append(run.font.size.pt)
                    if sizes:
                        size = max(sizes)
                    lines = para_lines(txt, size, max(colw[min(ci, len(colw)-1)] - 12, 10))
                    row_need = max(row_need, lines * size * 1.18 + 8)
                need += row_need
            if need > emu2pt(h) + 18:
                issues.append(f"S{idx}: TABLE overflow risk: need~{need:.0f}pt vs box {emu2pt(h):.0f}pt")
        elif shp.has_text_frame:
            tf = shp.text_frame
            txt = tf.text.strip()
            if not txt:
                continue
            wpt = emu2pt(w) - emu2pt(tf.margin_left or 0) - emu2pt(tf.margin_right or 0)
            req = frame_requirement(tf, max(wpt, 10))
            avail = emu2pt(h)
            if req > avail + 6:
                issues.append(f"S{idx}: TEXT overflow risk: need~{req:.0f}pt vs box {avail:.0f}pt :: '{txt[:38]}...'")

print(f"Slide size: {emu2pt(SW):.0f}x{emu2pt(SH):.0f} pt")
if issues:
    print(f"\n{len(issues)} POTENTIAL ISSUES:")
    for i in issues:
        print("  -", i)
else:
    print("\nNo overflow/out-of-bounds risks detected.")
