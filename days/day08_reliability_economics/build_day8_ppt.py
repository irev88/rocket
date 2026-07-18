# -*- coding: utf-8 -*-
"""
Build Day 8 PPT deck — 可靠性与经济性 (Reliability & Economics)
Bilingual CN/EN, consistent with Day 6/7 style.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

NAVY   = RGBColor(0x16, 0x32, 0x4F)
BLUE   = RGBColor(0x2E, 0x5E, 0x8C)
ORANGE = RGBColor(0xE0, 0x56, 0x1E)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GRAY   = RGBColor(0x5B, 0x66, 0x70)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF9)
CARD   = RGBColor(0xEE, 0xF3, 0xF8)
PEACH  = RGBColor(0xFB, 0xEE, 0xE6)
MINT   = RGBColor(0xE8, 0xF2, 0xEC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x33, 0x47, 0x5B)
FONT   = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
PAGE = [0]

def _set_font(run, size, bold=False, color=DARK, italic=False):
    f = run.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None: e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", FONT)

def rect(slide, x, y, w, h, fill, line=None, lw=1.0, round_=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try: shp.adjustments[0] = 0.05
        except: pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, space_after=2):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after)
        runs = ln if isinstance(ln, list) else [ln]
        for (t, size, bold, color) in runs:
            r = p.add_run(); r.text = t; _set_font(r, size, bold, color)
    return tb

def bullet_block(slide, x, y, w, h, items, size=14, color=DARK, gap=5, mark_color=ORANGE):
    lines = []
    for it in items:
        if isinstance(it, tuple):
            lines.append([("▪ ", size, True, mark_color), (it[0], size, True, NAVY), (it[1], size, False, color)])
        else:
            lines.append([("▪ ", size, True, mark_color), (it, size, False, color)])
    return text(slide, x, y, w, h, lines, space_after=gap)

def header(slide, title, sub=None):
    rect(slide, 0, 0, 13.333, 0.92, NAVY)
    rect(slide, 0, 0.92, 13.333, 0.045, ORANGE)
    text(slide, 0.55, 0.12, 11.5, 0.7, [(title, 25, True, WHITE)])
    if sub:
        text(slide, 9.2, 0.30, 3.6, 0.5, [(sub, 12, False, RGBColor(0xBF,0xD7,0xEA))], align=PP_ALIGN.RIGHT)
    PAGE[0] += 1
    rect(slide, 0, 7.18, 13.333, 0.32, LIGHT)
    text(slide, 0.55, 7.19, 9.5, 0.3,
         [("可重复使用运载火箭 AI 协同设计暑期课程 · Day 8 可靠性与经济性 · 2026-07-18", 10, False, GRAY)])
    text(slide, 12.2, 7.19, 0.7, 0.3, [(str(PAGE[0]), 10, False, GRAY)], align=PP_ALIGN.RIGHT)

def style_cell(cell, txt, size=12, bold=False, color=DARK, fill=None, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.MIDDLE):
    cell.margin_left = Pt(4); cell.margin_right = Pt(4)
    cell.margin_top = Pt(1.5); cell.margin_bottom = Pt(1.5)
    cell.vertical_anchor = anchor
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for i, (t, b, c) in enumerate(txt):
        r = p.add_run(); r.text = t; _set_font(r, size, b, c)

def make_table(slide, x, y, w, h, data, col_w=None, header_fill=NAVY, body_size=11,
               head_size=12, band=True, first_col_bold=False):
    rows, cols = len(data), len(data[0])
    gfx = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gfx.table
    tbl.first_row = False; tbl.horz_banding = False
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / total))
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                style_cell(cell, [(val, True, WHITE)], size=head_size, fill=header_fill, align=PP_ALIGN.CENTER)
            else:
                fill = WHITE
                if band and ri % 2 == 0: fill = CARD
                bold = first_col_bold and ci == 0
                col = NAVY if bold else DARK
                style_cell(cell, [(val, bold, col)], size=body_size, fill=fill)
    return tbl

def notes(slide, text_string):
    slide.notes_slide.notes_text_frame.text = text_string

# ============================================================ SLIDE 1 — TITLE
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 4.62, 13.333, 0.05, ORANGE)
text(s, 0.9, 1.05, 11.5, 0.5, [("可重复使用运载火箭 AI 协同设计暑期课程（11–20 July 2026）", 15, False, RGBColor(0xBF,0xD7,0xEA))])
text(s, 0.9, 1.65, 11.6, 1.6, [("可靠性与经济性分析", 42, True, WHITE),
                               ("（Reliability & Economics Analysis）", 22, False, RGBColor(0xE8,0xF2,0xEC))])
text(s, 0.9, 3.50, 11.5, 0.8, [("生命周期成本模型 · FMEA 失效分析 · 可靠性框图 · 过载缓解策略 · 双路径经济权衡", 16, False, RGBColor(0xBF,0xD7,0xEA))])
rect(s, 0.9, 4.95, 4.9, 1.5, RGBColor(0x1F,0x45,0x6E), round_=True)
text(s, 1.12, 5.12, 4.5, 1.2, [("Day 8 / 10 · 2026年7月18日", 16, True, WHITE),
                               ("交付物：成本与风险分析报告", 12.5, False, RGBColor(0xBF,0xD7,0xEA))])
rect(s, 6.0, 4.95, 6.4, 1.5, RGBColor(0x1F,0x45,0x6E), round_=True)
text(s, 6.22, 5.12, 6.1, 1.3, [("技术输入：Day 7 双闭合路径（Path A: 802t/20t; Path B: 600t/12t）", 12.5, False, RGBColor(0xBF,0xD7,0xEA)),
                               ("核心结论：Path A 单位成本最优 $1,883/kg；Path B 满足 $30M 限额", 12.5, False, RGBColor(0xBF,0xD7,0xEA))])
notes(s, "各位老师好，今天由我代表团队汇报第八天的核心交付物——可靠性与经济性分析。我们在Day 7两条物理闭合路径的基础上，建立了全生命周期成本模型，完成了FMEA失效模式分析，并提出了结构过载缓解策略。")

# ============================================================ SLIDE 2 — INPUTS FROM DAY 7
s = prs.slides.add_slide(BLANK); header(s, "技术输入：来自 Day 7 的双闭合基线", "Inputs from Day 7")
data = [
 ["参数", "Path A（载荷驱动）", "Path B（吨位驱动）"],
 ["起飞质量 GLOM", "801,600 kg (+34%)", "600,000 kg"],
 ["SSO 载荷", "20,000 kg", "12,000 kg"],
 ["S1 发动机", "12× Merlin 1D", "9× Merlin 1D"],
 ["S2 发动机", "4× MVac", "2× MVac"],
 ["S1 推进剂", "545 t", "384 t"],
 ["回收储备", "34.5 t", "34.5 t"],
 ["Max-Q", "33.1 kPa ✓", "31.5 kPa ✓"],
 ["SECO 峰值 g", "4.88 g ✓", "5.95 g → 节流至 4.85 g"],
]
make_table(s, 0.55, 1.25, 12.25, 4.8, data, col_w=[2.5,4.9,4.9], body_size=13, head_size=13.5)
notes(s, "Day 7 通过物理修复和多维迭代，确立了两条闭合路径。Path A 坚守 20t 载荷但 GLOM 增至 802t；Path B 坚守 600t 但载荷降至 12t。两者回收储备均锁定 34.5t。")

# ============================================================ SLIDE 3 — COST MODEL FORMULA
s = prs.slides.add_slide(BLANK); header(s, "参数化成本模型：公式体系", "Cost Model")
bullet_block(s, 0.55, 1.30, 6.0, 5.2, [
 ("单次发射边际成本：", "$C_{launch} = C_{amort} + C_{upper} + C_{fairing} + C_{refurb} + C_{ops} + C_{prop}$"),
 ("助推器摊薄（含5%保险系数）：", "$C_{amort} = C_{booster} × 1.05 / N$"),
 ("推进剂成本（O/F=2.56混合单价）：", "$C_{prop} = m_{prop} × $0.53/kg"),
 ("翻修分段函数：", "$C_{refurb}(1) = 0$ (首飞); $C_{refurb}(N>1) = $5.0M$"),
 ("海上回收运营：", "$C_{ops} = $1.5M/次 (固定)"),
], size=13, gap=8)

rect(s, 6.95, 1.30, 5.8, 5.2, CARD, round_=True)
text(s, 7.20, 1.55, 5.3, 4.7, [
 [("成本模型输入参数表:", 14, True, ORANGE)],
 [("", 8, False, DARK)],
 [("参数          Path A      Path B", 12, True, NAVY)],
 [("开发费        $2,200M    $2,000M", 11, False, DARK)],
 [("助推器造价    $31.2M     $24.0M", 11, False, DARK)],
 [("上级造价      $24.0M     $12.0M", 11, False, DARK)],
 [("整流罩造价    $4.6M      $4.0M", 11, False, DARK)],
 [("翻修费(N>1)   $5.0M      $5.0M", 11, False, DARK)],
 [("海回运营      $1.5M      $1.5M", 11, False, DARK)],
 [("推进剂总量    707 t      503 t", 11, False, DARK)],
 [("推进剂成本    $0.37M     $0.27M", 11, False, DARK)],
], space_after=2)
notes(s, "成本模型完整公式体系。翻修费为首飞0、后续飞$5M的分段函数。推进剂混合单价$0.53/kg由O/F=2.56加权得出。")

# ============================================================ SLIDE 4 — ECONOMIC TRADE TABLE
s = prs.slides.add_slide(BLANK); header(s, "经济权衡：双路径全寿命成本对比", "Economic Trade")
data_eco = [
 ["指标", "N", "Path A (802t)", "Path B (600t)", "抛弃型基线"],
 ["SSO 载荷", "—", "20,000 kg", "12,000 kg", "20,000 kg"],
 ["发射成本", "N=1", "$63.23M", "$42.97M", "$42.98M"],
 ["", "N=5", "$42.02M", "$27.81M", "—"],
 ["", "N=15", "$37.65M", "$24.45M", "—"],
 ["", "N=25", "$36.78M", "$23.78M", "—"],
 ["单位成本", "N=1", "$3,162/kg", "$3,581/kg", "$2,149/kg"],
 ["", "N=5", "$2,101/kg", "$2,318/kg", "—"],
 ["", "N=15", "$1,883/kg", "$2,038/kg", "—"],
 ["", "N=25", "$1,839/kg", "$1,982/kg", "—"],
]
make_table(s, 0.55, 1.25, 12.25, 5.2, data_eco, col_w=[1.5,0.8,2.5,2.5,2.5], body_size=12, head_size=12.5)
rect(s, 0.55, 6.55, 12.25, 0.5, MINT, round_=True)
text(s, 0.8, 6.58, 11.8, 0.45, [
 [("核心发现：", 13, True, GREEN),
  ("Path A 单位成本 $1,883/kg 比 Path B $2,038/kg 低 7.6%——规模经济效应；Path B 绝对成本 $24.45M 满足 <$30M 约束。", 12, False, DARK)],
])
notes(s, "经济权衡核心表。Path B满足绝对成本约束，但Path A单位成本更优。N=15是典型运营寿命点。")

# ============================================================ SLIDE 5 — RELIABILITY BLOCK DIAGRAM
s = prs.slides.add_slide(BLANK); header(s, "可靠性框图：任务成功与回收概率", "Reliability")
bullet_block(s, 0.55, 1.30, 5.5, 2.5, [
 ("任务成功概率（串联系联）：", "$P_{mission} = R_{S1} × R_{S2} = 0.992 × 0.990 = 98.2\\%$"),
 ("回收成功概率（5环节串联）：", "$P_{recovery} = P_{mission} × R_{entry} × R_{glide} × R_{burn} × R_{engage}$"),
 ("", "= 0.982 × 0.990 × 0.995 × 0.988 × 0.993 = 94.9%"),
 ("损耗率验证：", "5.1% 损耗 → 保险系数 $\\sigma_{ins}$ = 5% 完全合理"),
], size=13.5, gap=8)

# Reliability table
data_rel = [
 ["环节", "可靠度", "依据"],
 ["S1 上升 (9发冗余)", "99.2%", "发动机冗余；空中熄火容限"],
 ["S2 上升 (单发)", "99.0%", "无冗余；单点故障"],
 ["再入点火 (3发)", "99.0%", "三发重启；热防护"],
 ["栅格翼滑翔", "99.5%", "液压驱动器"],
 ["末段悬停猛击", "98.8%", "中心发动机重启；推进剂沉降"],
 ["缆索网捕获", "99.3%", "船端张紧网+液压阻尼"],
]
make_table(s, 6.85, 1.30, 5.9, 4.0, data_rel, col_w=[2.2,1.0,2.7], body_size=11.5, head_size=12)

rect(s, 0.55, 4.2, 5.5, 2.2, PEACH, round_=True)
text(s, 0.75, 4.35, 5.1, 1.9, [
 [("闭环节奏：", 13, True, NAVY)],
 [("任务成功 98.2% ← 客户载荷交付保障", 12, False, DARK)],
 [("回收成功 94.9% ← 助推器资产保障", 12, False, DARK)],
 [("每 20 次发射预计损失 1 枚助推器", 12, True, ORANGE)],
 [("完全融入成本模型的保险系数", 11, False, GRAY)],
], space_after=3)
notes(s, "可靠性框图：任务成功98.2%，回收成功94.9%。94.9%的损耗率验证了5%保险系数的合理性。")

# ============================================================ SLIDE 6 — FMEA TABLE (top 5)
s = prs.slides.add_slide(BLANK); header(s, "失效模式与影响分析 (FMEA) · 前5项", "FMEA")
data_fmea = [
 ["ID", "失效模式", "S", "O", "D", "RPN前", "缓解措施", "RPN后"],
 ["F-01", "末段发动机不重启", "9", "4", "5", "144", "冗余点火+避船弹道", "27"],
 ["F-02", "S2 过载超 5g", "8", "5", "1", "40", "40% 节流控制", "4"],
 ["F-03", "非对称挂接 (1-2耳)", "8", "4", "4", "128", "捕获环分散+磨损靴", "16"],
 ["F-06", "晃动致发动机断供", "8", "4", "3", "96", "挡板+防涡器+沉降烧", "8"],
 ["F-09", "RP-1 积碳裂纹", "6", "6", "4", "144", "复用护照+15飞寿命限", "16"],
]
make_table(s, 0.55, 1.25, 12.25, 3.5, data_fmea, col_w=[0.6,2.2,0.5,0.5,0.5,0.8,2.8,0.8], body_size=11, head_size=11.5, first_col_bold=True)

rect(s, 0.55, 4.95, 12.25, 1.6, LIGHT, round_=True)
bullet_block(s, 0.8, 5.05, 11.8, 1.4, [
 ("F-01 末段点火：", "置顶风险——朱雀三号与长征十二号A（2025年12月）均失利于此；冗余点火+避船弹道将 RPN 从 144 降至 27"),
 ("F-09 积碳：", "RP-1 特有老化机制；FBG 复用护照实现状态监控式检查，替代全面拆解；15飞寿命限"),
 ("完整 10 项 FMEA：", "见报告 §4 全表（F-01 至 F-10），所有 RPN 缓解后均 ≤27"),
], size=12, gap=4)
notes(s, "FMEA前5项。末段点火风险置顶（有2025年12月两次失利实证）；积碳为RP-1特有老化机制。完整10项见报告。")

# ============================================================ SLIDE 7 — OVER-ACCELERATION MITIGATION
s = prs.slides.add_slide(BLANK); header(s, "结构过载缓解：CR-D7-07 闭合", "Over-acceleration")
bullet_block(s, 0.55, 1.30, 5.8, 4.5, [
 ("问题：", "2× MVac 构型在 SECO 瞬间达到 6.2–6.4 g，超出 5.0 g 结构限制"),
 ("根因：", "推力/燃尽质量比过高（结构问题，非制导问题）"),
 ("三条缓解路径：", ""),
 ("  (A) 提高载荷过载容限至 6.5 g", "→ 拒绝：限制商业卫星市场"),
 ("  (B) 加固 S2 结构", "→ 拒绝：+350 kg 干重 → -800 kg 载荷"),
 ("  (C) S2 发动机节流至 40%", "→ ★ 选定：仅 -120 kg 载荷代价"),
], size=13, gap=6)

rect(s, 6.85, 1.30, 5.9, 4.5, MINT, round_=True)
text(s, 7.1, 1.55, 5.4, 4.0, [
 [("选定方案 (C) 技术细节:", 14, True, GREEN)],
 [("", 8, False, DARK)],
 [("• 当 S2 推进剂剩余 ≤10% 时（末 ~40 s）", 12.5, False, DARK)],
 [("  将 2× MVac 节流至 40%", 12.5, True, NAVY)],
 [("• 峰值加速度降至 4.85 g < 5.0 g ✓", 12.5, False, DARK)],
 [("• 重力损失增加极小", 12.5, False, DARK)],
 [("• 载荷代价仅 120 kg (<1% of 12 t)", 12.5, True, GREEN)],
 [("", 8, False, DARK)],
 [("Merlin 级发动机已有成熟", 12, False, GRAY)],
 [("40% 深度节流能力 (Day 3)", 12, False, GRAY)],
], space_after=2)
notes(s, "CR-D7-07闭合：选定发动机节流方案，代价仅120kg载荷。加固结构代价800kg，完全不可接受。")

# ============================================================ SLIDE 8 — FLEET CADENCE
s = prs.slides.add_slide(BLANK); header(s, "舰队规模与周转目标", "Fleet sizing")
bullet_block(s, 0.55, 1.30, 5.8, 2.8, [
 ("任务频次需求：", "≥ 10 次/年 (Day 1 约束)"),
 ("最短舰队规模：", "$N_{boosters} ≥ 10 × 30 / 365 ≈ 0.8 ⇒ 2–3 枚轮换$"),
 ("周转时间目标：", "早期 ≤30 天；成熟 ≤14 天"),
 ("设计寿命：", "15 次（冲刺 40 次）；会计折旧 25 次"),
], size=13.5, gap=7)

rect(s, 6.85, 1.30, 5.9, 2.8, CARD, round_=True)
data_fleet = [
 ["指标", "目标值", "外部锚点"],
 ["助推器寿命", "15 次 (冲刺 40)", "Block 5: 折旧 25/目标 40/纪录 36"],
 ["年发射频次", "≥ 10 次", "F9 现役: 165 次/年 (2025)"],
 ["周转时间", "≤ 30 天 (成熟 ≤14)", "现役均值 ~40 天"],
 ["翻修成本", "$5M/次 (保守)", "对照 Block 5: ~$0.3M"],
 ["活跃助推器", "2–3 枚", "10×30/365 ≈ 0.8 + 裕度"],
]
make_table(s, 6.95, 1.40, 5.7, 2.5, data_fleet, col_w=[1.5,1.8,2.4], body_size=11, head_size=11.5, first_col_bold=True)

rect(s, 0.55, 4.4, 12.25, 2.2, PEACH, round_=True)
text(s, 0.8, 4.5, 11.8, 0.35, [("成熟舰队经济 (C_refurb = $1.0M):", 14, True, NAVY)])
bullet_block(s, 0.8, 4.9, 11.8, 1.5, [
 ("Path A @ N=25:", "单位成本降至 $1,639/kg（较抛弃型基线 -24%）"),
 ("Path B @ N=25:", "单位成本降至 $1,648/kg（较抛弃型基线 -23%）"),
 ("两条路径均随翻修成熟实现显著经济优势", "——复用经济的核心驱动力是翻修成本从 $5M 降至 $1M 级"),
], size=12.5, gap=5)
notes(s, "舰队规模2-3枚轮换。成熟翻修$1M级时，两条路径均实现优于抛弃型的经济表现。")

# ============================================================ SLIDE 9 — DECISION SUMMARY
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 1.55, 13.333, 0.05, ORANGE)
text(s, 0.7, 0.45, 12.0, 0.9, [("Day 8 决策闭合：可靠性与经济性的系统级权衡", 26, True, WHITE)])
items = [
 ("成本模型闭合", "参数化公式体系（含分段翻修函数）完整可复现；Path A 单位成本 $1,883/kg 优于 Path B $2,038/kg（7.6%）"),
 ("可靠性闭合", "任务成功 98.2%、回收成功 94.9%；5% 损耗率完全融入保险系数；10 项 FMEA 全部 RPN ≤ 27"),
 ("结构过载闭合", "2× MVac 的 6.2 g 超限通过 40% 节流缓解至 4.85 g；代价仅 120 kg（<1% 载荷）"),
 ("运营概念闭合", "2–3 枚助推器轮换，早期 ≤30 天周转；成熟舰队（$1M 翻修）两条路径均优于抛弃型"),
 ("移交 Day 9", "两条闭合基线 + 全部量化风险/经济指标 → Day 9 系统集成技术评审 + Day 10 竞赛答辩"),
]
y = 1.95
for t_, d_ in items:
    rect(s, 0.7, y, 0.16, 0.80, ORANGE)
    text(s, 1.05, y, 3.0, 0.8, [(t_, 16, True, RGBColor(0xF2,0xA3,0x3C))])
    text(s, 3.9, y + 0.02, 8.9, 0.9, [(d_, 13.5, False, WHITE)])
    y += 1.02

rect(s, 0.7, 6.28, 12.0, 1.0, RGBColor(0x1F,0x45,0x6E), round_=True)
text(s, 0.95, 6.38, 11.6, 0.85, [
 [("本日技术审定意见：", 12, True, RGBColor(0xF2,0xA3,0x3C)),
  ("Day 8 成功完成了从物理闭合到工程-经济闭合的跨越。两条路径各有优势——Path A 适合追求单位成本最优的大批量星座客户，Path B 适合追求绝对成本可控的中小载荷客户。所有风险均已量化、所有指标均已闭合。", 11.5, False, RGBColor(0xBF,0xD7,0xEA))],
], space_after=3)
notes(s, "Day 8终审。成本、可靠性、结构、运营四维全部闭合。两条路径各有市场定位。移交Day 9系统集成和Day 10竞赛答辩。")

prs.save("presentation.pptx")
print("Saved presentation.pptx, slides:", len(prs.slides._sldIdLst))
